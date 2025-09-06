import pandas as pd 
import os 
import re 
import time 
from typing import List, Dict, Optional 
from transformers import AutoTokenizer 
from vllm import LLM, SamplingParams 
import math 

# 文档解析库 
try: 
    import fitz 
    from docx import Document 
    from pptx import Presentation 
except ImportError as e: 
    print(f"警告: 缺少文档解析库。请运行 'pip install PyMuPDF python-docx python-pptx' 安装。详细信息: {e}") 
    fitz, Document, Presentation = None, None, None 

class LLMDocumentSummarizer: 
    """ 
    使用vLLM加速，并利用多级MapReduceChain策略高效总结各种文档。 
    """ 
    def __init__(self, model_name_or_path: str, max_model_len: int = 40960, enable_thinking: bool = True): 
        print("正在加载模型和分词器...") 
        try: 
            self.tokenizer = AutoTokenizer.from_pretrained(model_name_or_path, trust_remote_code=True) 
            self.tokenizer.chat_template = ( 
                "{% for message in messages %}" 
                "{{ '<|im_start|>' + message['role'] + '\n' + message['content'] + '<|im_end|>' + '\n' }}" 
                "{% endfor %}" 
                "{% if add_generation_prompt %}" 
                "{{ '<|im_start|>assistant\n' }}" 
                "{% endif %}" 
            ) 
            
            self.llm = LLM( 
                model=model_name_or_path, 
                max_model_len=max_model_len, 
                trust_remote_code=True 
            ) 
            self.enable_thinking = enable_thinking 
            self.max_model_len = max_model_len 
            print(f"模型和分词器加载完成。Qwen3思考模式已{'启用' if self.enable_thinking else '禁用'}") 
        except Exception as e: 
            raise RuntimeError(f"错误：加载模型或分词器失败。请检查路径和环境配置。详细信息: {e}") 

    def _clean_response(self, text: str) -> str: 
        if text: 
            text = re.sub(r'<|im_end|>', '', text, flags=re.DOTALL) 
            text = text.strip() 
        return text 

    def _extract_text_from_file(self, file_path: str, sheet_name: Optional[str] = None) -> str: 
        file_extension = os.path.splitext(file_path)[1].lower() 
        content = "" 
        print(f"正在解析文件: {file_path}") 

        if file_extension in ['.xlsx', '.xls']: 
            try: 
                df = pd.read_excel(file_path, sheet_name=sheet_name) 
                content = df.to_markdown(index=False) 
            except Exception as e: 
                raise ValueError(f"无法读取Excel文件: {e}") 
        elif file_extension == '.txt': 
            try: 
                with open(file_path, 'r', encoding='utf-8') as f: 
                    content = f.read() 
            except Exception as e: 
                raise ValueError(f"无法读取TXT文件: {e}") 
        elif file_extension == '.pdf': 
            if fitz is None: 
                raise ImportError("PyMuPDF (fitz) 未安装。请安装: pip install PyMuPDF") 
            try: 
                with fitz.open(file_path) as doc: 
                    for page in doc: 
                        content += page.get_text() 
            except Exception as e: 
                raise ValueError(f"无法读取PDF文件: {e}") 
        elif file_extension == '.docx': 
            if Document is None: 
                raise ImportError("python-docx 未安装。请安装: pip install python-docx") 
            try: 
                doc = Document(file_path) 
                for para in doc.paragraphs: 
                    content += para.text + '\n' 
            except Exception as e: 
                raise ValueError(f"无法读取Word文件: {e}") 
        elif file_extension == '.pptx': 
            if Presentation is None: 
                raise ImportError("python-pptx 未安装。请安装: pip install python-pptx") 
            try: 
                prs = Presentation(file_path) 
                for slide in prs.slides: 
                    for shape in slide.shapes: 
                        if hasattr(shape, "text"): 
                            content += shape.text + '\n' 
            except Exception as e: 
                raise ValueError(f"无法读取PPT文件: {e}") 
        else: 
            raise ValueError(f"不支持的文件格式: {file_extension}") 
        
        content = re.sub(r'\s+', ' ', content).strip() 
        if not content: 
            raise ValueError("提取的文本内容为空。") 
            
        return content 

    def _recursive_summarize(self, texts: List[str], max_len: int, level: int = 1) -> str: 
        """ 
        递归地对文本块进行总结，直到汇总结果小于max_len。 
        """ 
        print(f"\n--- 步骤{level}: 开始第{level}级分块总结 ---") 
        prompts = [] 
        for text_chunk in texts: 
            messages = [ 
                { 
                    "role": "user", 
                    "content": f"请根据以下内容片段，生成一个简洁但全面的小结，以便后续进行最终汇总。\n\n内容片段：\n{text_chunk}\n\n小结：" 
                } 
            ] 
            prompt = self.tokenizer.apply_chat_template(messages, tokenize=False, add_generation_prompt=True, enable_thinking=self.enable_thinking) 
            prompts.append(prompt) 

        sub_summary_params = SamplingParams( 
            temperature=0.6, 
            top_p=0.95, 
            top_k=20, 
            max_tokens=self.max_model_len // 10, 
            n=1 
        ) 
        
        try: 
            outputs = self.llm.generate(prompts, sub_summary_params) 
            sub_summaries = [self._clean_response(output.outputs[0].text) for output in outputs] 
            combined_summary = "\n\n---\n\n".join(sub_summaries) 
            
            summary_tokens = self.tokenizer.encode(combined_summary) 
            if len(summary_tokens) > max_len: 
                print(f"第{level}级汇总后文本长度为 {len(summary_tokens)}，超过限制 {max_len}。进行下一级汇总。") 
                
                # 计算新的分块大小 
                new_chunk_size = math.ceil(len(summary_tokens) / (len(summary_tokens) // max_len + 1)) 
                new_chunks = [summary_tokens[i:i + new_chunk_size] for i in range(0, len(summary_tokens), new_chunk_size)] 
                new_text_chunks = [self.tokenizer.decode(c, skip_special_tokens=True) for c in new_chunks] 

                return self._recursive_summarize(new_text_chunks, max_len, level + 1) 
            else: 
                return combined_summary 
        except Exception as e: 
            raise RuntimeError(f"在第 {level} 级汇总时调用模型发生错误: {e}") 

    def summarize_document(self, file_path: str, question_text: str, chunk_token_size: int = 2000, sheet_name: Optional[str] = None) -> str: 
        if not os.path.exists(file_path): 
            return f"错误：未找到文件 '{file_path}'。" 

        print(f"\n--- 步骤1: 提取并分块文档文本 ---") 
        try: 
            full_text = self._extract_text_from_file(file_path, sheet_name) 
            print(f"成功提取文件内容，总字符数: {len(full_text)}。") 
        except (ValueError, ImportError) as e: 
            return str(e) 
        
        full_tokens = self.tokenizer.encode(full_text) 
        text_chunks = [self.tokenizer.decode(full_tokens[i:i + chunk_token_size], skip_special_tokens=True) for i in range(0, len(full_tokens), chunk_token_size)] 
        print(f"文本已分割成 {len(text_chunks)} 个块，每块大小约为 {chunk_token_size} 个Token。") 

        try: 
            all_sub_summaries_text = self._recursive_summarize(text_chunks, max_len=self.max_model_len) 
            
            print("\n--- 步骤2: 汇总所有小结，生成最终总结 ---") 
            final_summary_instruction = f"我已将一个大型文档分块总结，以下是所有的分块小结内容。请根据这些小结，生成一个详细且连贯的最终总结，以回答我最初的问题：'{question_text}'。\n\n所有小结：\n{all_sub_summaries_text}" 
            
            final_prompt = self.tokenizer.apply_chat_template([{"role": "user", "content": final_summary_instruction}], tokenize=False, add_generation_prompt=True, enable_thinking=self.enable_thinking) 
            
            final_summary_params = SamplingParams( 
                temperature=0.7, 
                top_p=0.95, 
                top_k=20, 
                max_tokens=self.max_model_len, 
            ) 

            start_time = time.time() 
            outputs = self.llm.generate([final_prompt], final_summary_params) 
            end_time = time.time() 
            
            final_summary = self._clean_response(outputs[0].outputs[0].text) 
            print(f"最终总结用时: {end_time - start_time:.2f} 秒") 
            return final_summary 
        except RuntimeError as e: 
            return str(e) 


# --- 脚本主入口 --- 
if __name__ == "__main__": 
    model_path = r"/DATA/f60055380/Qwen3_8B_TL/Model/qwen3-8b-hf" 
    
    # 示例: 你可以修改这里的文件路径和问题 
    file_to_summarize = r"/DATA/f60055380/Qwen3_8B_TL/data/xlsx/问题集1200.xlsx" 
    question = "总结不同维度不同分数的规律。" 

    print("="*50) 
    print(f"任务开始: 总结文档内容") 
    print("="*50) 

    try: 
        summarizer = LLMDocumentSummarizer(model_path, enable_thinking=True)    # 思考模式关关闭填 False 
        
        summary_result = summarizer.summarize_document( 
            file_to_summarize, 
            question, 
            chunk_token_size=2500,    # 每次token数，小于4k效果较好 
            sheet_name="Sheet1" if file_to_summarize.endswith(('.xlsx', '.xls')) else None 
        ) 
        
        print("\n\n" + "="*20 + " 最终总结结果 " + "="*20) 
        print(summary_result) 
        print("="*58) 
        
    except RuntimeError as e: 
        print(f"\n处理过程中发生致命错误: {e}")

