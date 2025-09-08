# -*- coding: utf-8 -*- 

import os 
import re 
import time 
import gc 
from concurrent.futures import ThreadPoolExecutor, as_completed 
from typing import List, Optional 

# --- 核心依赖 --- 
import pandas as pd 
import numpy as np 
import torch 
from transformers import AutoTokenizer, AutoModel 
from vllm import LLM, SamplingParams 
from faiss import IndexFlatL2 

# --- FastAPI 和 Pydantic 组件 --- 
from fastapi import FastAPI, HTTPException, File, UploadFile, Form 
from fastapi.middleware.cors import CORSMiddleware 
from pydantic import BaseModel 

# --- LangChain 组件 --- 
from langchain_text_splitters import RecursiveCharacterTextSplitter 

# --- 文档解析库 (通用) --- 
try: 
    import fitz  # PyMuPDF 
    from docx import Document 
    from pptx import Presentation 
except ImportError as e: 
    print(f"警告: 缺少文档解析库。详细信息: {e}") 
    fitz, Document, Presentation = None, None, None 


# --- 将你的类和辅助函数放在这里 --- 

# 获取可用GPU数量 
def get_num_gpus_from_env(): 
    cuda_visible = os.environ.get('CUDA_VISIBLE_DEVICES', None) 
    if cuda_visible: 
        return len([x for x in cuda_visible.split(',') if x.strip().isdigit()]) 
    try: 
        import torch 
        return torch.cuda.device_count() 
    except Exception: 
        return 1 


class AdvancedRAGSummarizer: 
    """ 
    一个集成了RAG、多级MapReduce和模型特定功能的高级文档处理类。 
    (已优化显存管理和数据类型处理) 
    """ 

    # 保留你原始代码中的所有方法，只在需要的地方稍作调整 
    # ... (此处省略，请将你提供的类代码完整复制到此处) ... 
    def __init__(self, model_name_or_path: str, embedding_model_name: str, max_model_len: int = 40960, enable_thinking: bool = True): 
        """ 
        初始化配置，但不立即加载模型以节省显存。 
        """ 
        self.llm_model_path = model_name_or_path 
        self.embedding_model_path = embedding_model_name 
        self.max_model_len = max_model_len 
        self.enable_thinking = enable_thinking 

        # 将模型和分词器初始化为 None 
        self.llm = None 
        self.tokenizer = None 
        self.embedding_model = None 
        self.embedding_tokenizer = None 
        
        # 预先确定数据类型 
        if torch.cuda.is_available(): 
            self.dtype = torch.bfloat16 if torch.cuda.is_bf16_supported() else torch.float16 
        else: 
            self.dtype = torch.float32 

        print("="*50) 
        print("AdvancedRAGSummarizer 初始化完成。") 
        print(f"LLM 模型路径: {self.llm_model_path}") 
        print(f"Embedding 模型路径: {self.embedding_model_path}") 
        print(f"将使用 {self.dtype} 数据类型加载模型。") 
        print("="*50) 

    def _load_llm(self): 
        """按需加载LLM模型。""" 
        if self.llm is None: 
            print("\n--- 正在加载LLM模型和分词器... ---") 
            try: 
                num_gpus = get_num_gpus_from_env() 
                self.tokenizer = AutoTokenizer.from_pretrained(self.llm_model_path, trust_remote_code=True) 
                self.tokenizer.chat_template = ( 
                    "{% for message in messages %}" 
                    "{{ '<|im_start|>' + message['role'] + '\n' + message['content'] + '<|im_end|>' + '\n' }}" 
                    "{% endfor %}" 
                    "{% if add_generation_prompt %}" 
                    "{{ '<|im_start|>assistant\n' }}" 
                    "{% endif %}" 
                ) 
                self.llm = LLM( 
                    model=self.llm_model_path, 
                    max_model_len=self.max_model_len, 
                    trust_remote_code=True, 
                    tensor_parallel_size=num_gpus if num_gpus > 1 else 1, 
                    dtype=self.dtype 
                ) 
                print("--- LLM模型和分词器加载完成。 ---") 
            except Exception as e: 
                raise RuntimeError(f"错误：加载LLM失败。详细信息: {e}") 

    def _unload_llm(self): 
        """卸载LLM模型并清理显存。""" 
        if self.llm is not None: 
            print("\n--- 正在卸载LLM模型... ---") 
            del self.llm 
            del self.tokenizer 
            self.llm = None 
            self.tokenizer = None 
            gc.collect() 
            torch.cuda.empty_cache() 
            print("--- LLM模型已卸载并清理显存。 ---") 

    def _load_embedding_model(self): 
        """按需加载Embedding模型。""" 
        if self.embedding_model is None: 
            print("\n--- 正在加载Embedding模型... ---") 
            try: 
                self.embedding_tokenizer = AutoTokenizer.from_pretrained(self.embedding_model_path, trust_remote_code=True) 
                self.embedding_model = AutoModel.from_pretrained( 
                    self.embedding_model_path, 
                    trust_remote_code=True, 
                    device_map="auto", 
                    torch_dtype=self.dtype 
                ) 
                self.embedding_model.eval() 
                print("--- Embedding模型加载完成。 ---") 
            except Exception as e: 
                raise RuntimeError(f"错误：加载Embedding模型失败。详细信息: {e}") 

    def _unload_embedding_model(self): 
        """卸载Embedding模型并清理显存。""" 
        if self.embedding_model is not None: 
            print("\n--- 正在卸载Embedding模型... ---") 
            del self.embedding_model 
            del self.embedding_tokenizer 
            self.embedding_model = None 
            self.embedding_tokenizer = None 
            gc.collect() 
            torch.cuda.empty_cache() 
            print("--- Embedding模型已卸载并清理显存。 ---") 

    def _clean_response(self, text: str) -> str: 
        if text: 
            text = re.sub(r'<\|im_end\|>', '', text, flags=re.DOTALL) 
            text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL) 
            text = text.strip() 
        return text 

    def _get_sampling_params(self, is_thinking: bool): 
        if is_thinking: 
            return SamplingParams(temperature=0.6, top_p=0.95, top_k=20, max_tokens=2048) 
        else: 
            return SamplingParams(temperature=0.7, top_p=0.8, top_k=20, max_tokens=2048) 

    def _recursive_summarize(self, texts: List[str], level: int = 1) -> str: 
        self._load_llm() # 确保LLM已加载 
        print(f"\n--- 第 {level} 级分块总结 (共 {len(texts)} 块) ---") 
        prompts = [] 
        soft_switch = "/think" if self.enable_thinking else "/no_think" 
        for text_chunk in texts: 
            messages = [{"role": "user", "content": f"{soft_switch} 请根据以下内容片段，生成一个简洁但全面的小结，以便后续进行最终汇总。\n\n内容片段：\n{text_chunk}\n\n小结："}] 
            prompt = self.tokenizer.apply_chat_template(messages, tokenize=False, add_generation_prompt=True) 
            prompts.append(prompt) 
        
        sub_summary_params = self._get_sampling_params(self.enable_thinking) 
        sub_summary_params.max_tokens = max(512, self.max_model_len // 10) 
        
        try: 
            outputs = self.llm.generate(prompts, sub_summary_params) 
            sub_summaries = [self._clean_response(output.outputs[0].text) for output in outputs] 
            
            combined_summary = "\n\n---\n\n".join(sub_summaries) 
            summary_tokens = self.tokenizer.encode(combined_summary) 
            
            if len(summary_tokens) > self.max_model_len * 0.8: 
                print(f"第{level}级汇总后Token数 {len(summary_tokens)} 超过限制，进行下一级汇总。") 
                splitter = RecursiveCharacterTextSplitter( 
                    chunk_size=self.max_model_len, 
                    chunk_overlap=200, 
                    length_function=lambda x: len(self.tokenizer.encode(x)) 
                ) 
                new_chunks = splitter.split_text(combined_summary) 
                return self._recursive_summarize(new_chunks, level + 1) 
            else: 
                return combined_summary 
        except Exception as e: 
            raise RuntimeError(f"在第 {level} 级汇总时调用模型发生错误: {e}") 

    def _generate_embeddings(self, texts: List[str], batch_size: int = 16) -> List[List[float]]: 
        self._load_embedding_model() # 确保Embedding模型已加载 
        all_embeddings = [] 
        for i in range(0, len(texts), batch_size): 
            batch_texts = texts[i:i + batch_size] 
            inputs = self.embedding_tokenizer(batch_texts, padding=True, truncation=True, return_tensors='pt').to(self.embedding_model.device) 
            with torch.no_grad(): 
                outputs = self.embedding_model(**inputs) 
            
            embeddings = outputs.last_hidden_state[:, 0, :] 
            embeddings = torch.nn.functional.normalize(embeddings, p=2, dim=1) 
            
            # --- 关键修复：在这里将数据类型转换为 float32 --- 
            all_embeddings.extend(embeddings.to(torch.float32).cpu().numpy().tolist()) 
            
            # 主动清理 
            del inputs, outputs, embeddings 
        
        torch.cuda.empty_cache() 
        return all_embeddings 

    def process_document(self, file_path: str, question_text: str, chunk_size_chars: int = 1000, sheet_name: Optional[str] = None) -> str: 
        if not os.path.exists(file_path): 
            return f"错误：未找到文件 '{file_path}'。" 

        try: 
            # --- 步骤 1: 文本分块 --- 
            print("\n--- 步骤1: 提取和分块文本 ---") 
            text_chunks = self._chunk_file_content(file_path, chunk_size_chars, sheet_name) 

            # --- 步骤 2: RAG检索 (仅使用Embedding模型) --- 
            print("\n--- 步骤2: 构建向量索引并进行RAG检索 ---") 
            start_time = time.time() 
            
            # 按需加载Embedding模型 
            self._load_embedding_model() 
            
            print("正在生成文本块嵌入向量...") 
            chunk_vectors = self._generate_embeddings(text_chunks) 
            
            question_instructed = f"Given a user's question, retrieve relevant passages that answer the question.\n{question_text}" 
            question_vector = self._generate_embeddings([question_instructed])[0] 
            
            # **关键：用完后立即卸载Embedding模型，为LLM腾出空间** 
            self._unload_embedding_model() 

            d = len(chunk_vectors[0]) 
            chunk_vectors_np = np.array(chunk_vectors, dtype=np.float32) 
            index = IndexFlatL2(d) 
            index.add(chunk_vectors_np) 
            print(f"向量索引构建完成，用时: {time.time() - start_time:.2f} 秒") 

            print(f"正在根据问题检索最相关的文本块...") 
            query_vector_np = np.array([question_vector], dtype=np.float32) 
            D, I = index.search(query_vector_np, k=min(10, len(text_chunks))) 
            retrieved_chunks = [text_chunks[i] for i in I[0]] 
            print(f"成功检索到 {len(retrieved_chunks)} 个相关文本块。") 
            
            # 清理向量数据 
            del chunk_vectors, question_vector, chunk_vectors_np, index, query_vector_np 
            gc.collect() 
            torch.cuda.empty_cache() 


            # --- 步骤 3 & 4: 总结与回答 (仅使用LLM) --- 
            # 按需加载LLM 
            self._load_llm() 

            all_sub_summaries_text = self._recursive_summarize(retrieved_chunks) 

            print("\n--- 步骤3: 汇总所有小结，生成最终总结 ---") 
            final_summary_instruction = f"我已从一个大型文档中检索出与问题相关的片段，并对它们进行了初步总结。以下是所有相关的小结内容，请根据这些内容，生成一个详细且连贯的最终总结。\n\n所有小结：\n{all_sub_summaries_text}" 
            final_summary_instruction = f"{'/think' if self.enable_thinking else '/no_think'} " + final_summary_instruction 
            final_prompt = self.tokenizer.apply_chat_template([{"role": "user", "content": final_summary_instruction}], tokenize=False, add_generation_prompt=True) 
            final_summary_params = self._get_sampling_params(self.enable_thinking) 
            outputs = self.llm.generate([final_prompt], final_summary_params) 
            final_summary = self._clean_response(outputs[0].outputs[0].text) 

            print("\n--- 步骤4: 基于最终总结，回答用户最初的问题 ---") 
            answer_instruction = f"现在，请根据以下这份总结报告，清晰、准确地回答我最初的问题。\n\n总结报告：\n{final_summary}\n\n我的问题是：'{question_text}'\n\n请回答：" 
            answer_instruction = f"{'/think' if self.enable_thinking else '/no_think'} " + answer_instruction 
            answer_prompt = self.tokenizer.apply_chat_template([{"role": "user", "content": answer_instruction}], tokenize=False, add_generation_prompt=True) 
            answer_params = self._get_sampling_params(self.enable_thinking) 
            outputs = self.llm.generate([answer_prompt], answer_params) 
            final_answer = self._clean_response(outputs[0].outputs[0].text) 

            return f"✅ **处理完成！**\n\n" + "="*20 + " **综合总结** " + "="*20 + f"\n\n{final_summary}\n\n" + "="*20 + " **问题回答** " + "="*20 + f"\n\n{final_answer}" 

        except (ValueError, ImportError, RuntimeError) as e: 
            raise HTTPException(status_code=500, detail=f"处理过程中发生错误: {e}") 
        finally: 
            # **关键：确保任务结束时卸载所有模型** 
            self._unload_llm() 
            self._unload_embedding_model() 

    def _chunk_file_content(self, file_path, chunk_size_chars, sheet_name): 
        """将文件解析和分块逻辑提取为一个独立方法。""" 
        if file_path.endswith(('.xlsx', '.xls', '.csv')): 
            if file_path.endswith('.csv'): 
                df = pd.read_csv(file_path) 
            else: 
                df = pd.read_excel(file_path, sheet_name=sheet_name) if sheet_name else pd.read_excel(file_path) 
            row_blocks = [] 
            for idx, row in df.iterrows(): 
                block = [f"{col}: {row[col]}" for col in df.columns] 
                row_blocks.append("\n".join(block)) 
            
            text_chunks = [] 
            cur_chunk = "" 
            for block in row_blocks: 
                if not cur_chunk or len(cur_chunk) + len(block) < chunk_size_chars: 
                    cur_chunk += ("\n---\n" + block) if cur_chunk else block 
                else: 
                    text_chunks.append(cur_chunk) 
                    cur_chunk = block 
            if cur_chunk: 
                text_chunks.append(cur_chunk) 
            print(f"表格已结构化分块为 {len(text_chunks)} 个块。") 
        else: 
            full_text = self._extract_text_from_file(file_path, sheet_name) 
            splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size_chars, chunk_overlap=100, length_function=len) 
            text_chunks = splitter.split_text(full_text) 
            print(f"文本已分割成 {len(text_chunks)} 个语义块。") 
        return text_chunks 

    def _extract_text_from_file(self, file_path: str, sheet_name: Optional[str] = None) -> str: 
        ext = os.path.splitext(file_path)[1].lower() 
        full_text = "" 
        if ext == '.pdf': 
            if not fitz: raise ImportError("未安装 PyMuPDF (fitz)。请运行 `pip install PyMuPDF`。") 
            with fitz.open(file_path) as doc: 
                for page in doc: 
                    full_text += page.get_text() 
        elif ext == '.docx': 
            if not Document: raise ImportError("未安装 python-docx。请运行 `pip install python-docx`。") 
            doc = Document(file_path) 
            full_text = "\n".join([p.text for p in doc.paragraphs]) 
        elif ext == '.pptx': 
            if not Presentation: raise ImportError("未安装 python-pptx。请运行 `pip install python-pptx`。") 
            prs = Presentation(file_path) 
            for slide in prs.slides: 
                for shape in slide.shapes: 
                    if hasattr(shape, "text_frame") and shape.text_frame: 
                        full_text += shape.text_frame.text + "\n" 
        elif ext in ['.txt', '.md', '.log']: 
            with open(file_path, 'r', encoding='utf-8') as f: 
                full_text = f.read() 
        elif ext in ['.xlsx', '.xls', '.csv']: 
            pass 
        else: 
            raise ValueError(f"不支持的文件类型: {ext}") 
        return full_text 


# --- FastAPI 接口部分 --- 

# 1. 实例化 FastAPI 应用 
app = FastAPI( 
    title="Advanced RAG Summarizer API", 
    description="一个用于文档RAG检索和总结的高级API接口。", 
    version="1.0.0", 
) 
# 2. 配置 CORS 中间件 
app.add_middleware( 
    CORSMiddleware, 
    allow_origins=["*"], 
    allow_credentials=True, 
    allow_methods=["*"], 
    allow_headers=["*"], 
) 

# 3. 移除旧的 DocumentRequest Pydantic 模型，因为现在请求参数是多部分表单 
# class DocumentRequest(BaseModel): 
#     ... 

# 4. 重新定义API接口，使用 @app.post 装饰器 
@app.post("/summarize-document/") 
async def summarize_document( 
    file: UploadFile = File(...), # 接收上传的文件 
    question: str = Form(...), # 使用 Form 接收字符串参数 
    llm_model_path: str = Form(...), 
    embedding_model_path: str = Form(...), 
    max_model_len: int = Form(40960), 
    enable_thinking: bool = Form(True), 
    chunk_size_chars: int = Form(1000), 
    sheet_name: Optional[str] = Form(None) 
): 
    """ 
    **API接口：文档上传与 RAG 总结** 
    
    该接口接收一个上传的文件和问题，利用 RAG 技术对文档进行总结并回答问题。 
    """ 
    temp_file_path = None 
    try: 
        # 在服务器上创建一个临时文件路径来保存上传的文件 
        temp_dir = "temp_uploads" 
        os.makedirs(temp_dir, exist_ok=True) 
        temp_file_path = os.path.join(temp_dir, file.filename) 
        
        # 将上传的文件内容写入临时文件 
        with open(temp_file_path, "wb") as buffer: 
            contents = await file.read() # 使用 await 读取文件内容 
            buffer.write(contents) 

        # 实例化你的RAG摘要器 
        summarizer = AdvancedRAGSummarizer( 
            model_name_or_path=llm_model_path, 
            embedding_model_name=embedding_model_path, 
            max_model_len=max_model_len, 
            enable_thinking=enable_thinking 
        ) 

        # 调用你的核心处理方法，传入临时文件路径 
        result = summarizer.process_document( 
            temp_file_path, 
            question, 
            chunk_size_chars, 
            sheet_name 
        ) 

        return {"message": "处理成功", "result": result} 
        
    except HTTPException as e: 
        raise e 
    except Exception as e: 
        raise HTTPException(status_code=500, detail=f"内部服务器错误: {str(e)}") 
    finally: 
        # 任务完成后，记得删除临时文件 
        if temp_file_path and os.path.exists(temp_file_path): 
            os.remove(temp_file_path) 

# ... (保留原有的 get_num_gpus_from_env 和 AdvancedRAGSummarizer 类) ... 

if __name__ == "__main__": 
    import uvicorn 
    uvicorn.run(app, host="0.0.0.0", port=8000)

