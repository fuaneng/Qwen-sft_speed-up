# -*- coding: utf-8 -*-

import os
import re
import time
from typing import List, Optional

# --- 核心依赖 ---
import pandas as pd
from transformers import AutoTokenizer
from vllm import LLM, SamplingParams

# --- LangChain 组件 (版本2和3的优点) ---
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.faiss import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings # 使用本地开源Embedding模型

# --- 文档解析库 (通用) ---
try:
    import fitz  # PyMuPDF
    from docx import Document
    from pptx import Presentation
except ImportError as e:
    print(f"警告: 缺少文档解析库。请运行 'pip install PyMuPDF python-docx python-pptx' 安装。详细信息: {e}")
    fitz, Document, Presentation = None, None, None


class AdvancedRAGSummarizer:
    """
    一个集成了RAG、多级MapReduce和模型特定功能（如Qwen3的“思考模式”）的高级文档处理类。
    - 版本3的RAG架构: 用于高效检索相关信息。
    - 版本2的智能文本分割和解耦流程: 用于保证处理逻辑和输出质量。
    - 版本1的模型特定优化: 用于最大化模型性能。
    """

    def __init__(self, model_name_or_path: str, embedding_model_name: str = "BAAI/bge-large-zh-v1.5", max_model_len: int = 40960, enable_thinking: bool = True):
        """
        初始化模型、分词器和RAG组件。
        """
        print("="*50)
        print("正在加载LLM模型和分词器...")
        try:
            # 恢复版本1的“思考模式”逻辑
            self.tokenizer = AutoTokenizer.from_pretrained(model_name_or_path, trust_remote_code=True)
            self.tokenizer.chat_template = (
                "{% for message in messages %}"
                "{{ '<|im_start|>' + message['role'] + '\n' + message['content'] + '<|im_end|>' + '\n' }}"
                "{% endfor %}"
                "{% if add_generation_prompt %}"
                "{{ '<|im_start|>assistant\n' }}"
                "{% endif %}"
            )
            
            self.llm = LLM(
                model=model_name_or_path,
                max_model_len=max_model_len,
                trust_remote_code=True
            )
            self.enable_thinking = enable_thinking
            self.max_model_len = max_model_len
            print(f"LLM模型和分词器加载完成。Qwen3思考模式已{'启用' if self.enable_thinking else '禁用'}")

            # 初始化RAG组件，使用开源本地Embedding模型
            print(f"正在加载Embedding模型: {embedding_model_name}...")
            self.embeddings = HuggingFaceEmbeddings(model_name=embedding_model_name)
            self.vector_store = None
            print("Embedding模型加载完成。")
            print("="*50)

        except Exception as e:
            raise RuntimeError(f"错误：加载模型失败。请检查路径和环境配置。详细信息: {e}")

    def _clean_response(self, text: str) -> str:
        """清理模型输出中的特殊标记。"""
        if text:
            text = re.sub(r'<|im_end|>', '', text, flags=re.DOTALL)
            text = text.strip()
        return text

    def _extract_text_from_file(self, file_path: str, sheet_name: Optional[str] = None) -> str:
        """
        从各种格式的文件中提取纯文本内容。
        (采用了版本2的优化写法)
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        content = ""
        print(f"正在解析文件: {file_path}...")

        try:
            if file_extension in ['.xlsx', '.xls']:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content = df.to_markdown(index=False)
            elif file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif file_extension == '.pdf':
                if fitz is None: raise ImportError("PyMuPDF (fitz) 未安装。")
                with fitz.open(file_path) as doc:
                    content = "\n".join(page.get_text() for page in doc)
            elif file_extension == '.docx':
                if Document is None: raise ImportError("python-docx 未安装。")
                doc = Document(file_path)
                content = "\n".join(para.text for para in doc.paragraphs)
            elif file_extension == '.pptx':
                if Presentation is None: raise ImportError("python-pptx 未安装。")
                prs = Presentation(file_path)
                content = "\n".join(
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes if hasattr(shape, "text")
                )
            else:
                raise ValueError(f"不支持的文件格式: {file_extension}")
        except Exception as e:
            raise ValueError(f"无法读取文件 '{file_path}': {e}")

        content = re.sub(r'\s+', ' ', content).strip()
        if not content:
            raise ValueError("提取的文本内容为空。")
        
        print(f"文件解析成功，总字符数: {len(content)}")
        return content

    def _recursive_summarize(self, texts: List[str], level: int = 1) -> str:
        """
        递归地对文本块进行总结，直到汇总结果小于模型长度限制。
        (采用了版本2的LangChain分割器进行递归)
        """
        print(f"\n--- 第 {level} 级分块总结 (共 {len(texts)} 块) ---")
        prompts = []
        for text_chunk in texts:
            messages = [
                {"role": "user", "content": f"请根据以下内容片段，生成一个简洁但全面的小结，以便后续进行最终汇总。\n\n内容片段：\n{text_chunk}\n\n小结："}
            ]
            # 恢复版本1的“思考模式”调用
            prompt = self.tokenizer.apply_chat_template(messages, tokenize=False, add_generation_prompt=True, enable_thinking=self.enable_thinking)
            prompts.append(prompt)

        # 恢复版本1的动态max_tokens计算
        sub_summary_params = SamplingParams(
            temperature=0.6, top_p=0.95, top_k=20, 
            max_tokens=max(512, self.max_model_len // 10), # 取一个合理的值
            n=1
        )
        
        try:
            outputs = self.llm.generate(prompts, sub_summary_params)
            sub_summaries = [self._clean_response(output.outputs[0].text) for output in outputs]
            combined_summary = "\n\n---\n\n".join(sub_summaries)
            
            summary_tokens = self.tokenizer.encode(combined_summary)
            if len(summary_tokens) > self.max_model_len * 0.8: # 使用80%作为递归阈值
                print(f"第{level}级汇总后Token数 {len(summary_tokens)} 超过限制，进行下一级汇总。")
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=self.max_model_len, 
                    chunk_overlap=200, 
                    length_function=lambda x: len(self.tokenizer.encode(x))
                )
                new_chunks = splitter.split_text(combined_summary)
                return self._recursive_summarize(new_chunks, level + 1)
            else:
                return combined_summary
        except Exception as e:
            raise RuntimeError(f"在第 {level} 级汇总时调用模型发生错误: {e}")

    def process_document(self, file_path: str, question_text: str, chunk_size_chars: int = 1000, sheet_name: Optional[str] = None) -> str:
        """
        完整处理流程：提取 -> 分块 -> RAG检索 -> 总结 -> 回答。
        """
        if not os.path.exists(file_path):
            return f"错误：未找到文件 '{file_path}'。"

        try:
            # --- 步骤 1: 提取与分块 (版本2的逻辑) ---
            print("\n--- 步骤1: 提取并使用LangChain分割文本 ---")
            full_text = self._extract_text_from_file(file_path, sheet_name)
            splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size_chars, chunk_overlap=100, length_function=len)
            text_chunks = splitter.split_text(full_text)
            print(f"文本已分割成 {len(text_chunks)} 个语义块。")

            # --- 步骤 2: 构建向量索引并进行RAG检索 (版本3的逻辑) ---
            print("\n--- 步骤2: 构建向量索引并进行RAG检索 ---")
            start_time = time.time()
            self.vector_store = FAISS.from_texts(text_chunks, self.embeddings)
            print(f"向量索引构建完成，用时: {time.time() - start_time:.2f} 秒")

            print(f"正在根据问题检索最相关的文本块...")
            relevant_docs = self.vector_store.similarity_search(question_text, k=min(10, len(text_chunks))) # 检索最多10个相关块
            retrieved_chunks = [doc.page_content for doc in relevant_docs]
            print(f"成功检索到 {len(retrieved_chunks)} 个相关文本块。")

            # --- 步骤 3: 对检索出的内容进行分级总结 ---
            all_sub_summaries_text = self._recursive_summarize(retrieved_chunks)
            
            # --- 步骤 4: 汇总生成通用总结 (版本2的解耦逻辑) ---
            print("\n--- 步骤3: 汇总所有小结，生成最终总结 ---")
            final_summary_instruction = f"我已从一个大型文档中检索出与问题相关的片段，并对它们进行了初步总结。以下是所有相关的小结内容，请根据这些内容，生成一个详细且连贯的最终总结。\n\n所有小结：\n{all_sub_summaries_text}"
            
            final_prompt = self.tokenizer.apply_chat_template([{"role": "user", "content": final_summary_instruction}], tokenize=False, add_generation_prompt=True, enable_thinking=self.enable_thinking)
            
            final_summary_params = SamplingParams(
                temperature=0.7, top_p=0.95, top_k=20, max_tokens=2048
            )

            outputs = self.llm.generate([final_prompt], final_summary_params)
            final_summary = self._clean_response(outputs[0].outputs[0].text)

            # --- 步骤 5: 基于通用总结回答具体问题 (版本2的解耦逻辑) ---
            print("\n--- 步骤4: 基于最终总结，回答用户最初的问题 ---")
            answer_instruction = f"现在，请根据以下这份总结报告，清晰、准确地回答我最初的问题。\n\n总结报告：\n{final_summary}\n\n我的问题是：'{question_text}'\n\n请回答："
            
            answer_prompt = self.tokenizer.apply_chat_template([{"role": "user", "content": answer_instruction}], tokenize=False, add_generation_prompt=True, enable_thinking=self.enable_thinking)
            
            answer_params = SamplingParams(
                temperature=0.1, top_p=0.9, top_k=10, max_tokens=1024
            )
            
            outputs = self.llm.generate([answer_prompt], answer_params)
            final_answer = self._clean_response(outputs[0].outputs[0].text)

            return f"✅ **处理完成！**\n\n" + "="*20 + " **综合总结** " + "="*20 + f"\n\n{final_summary}\n\n" + "="*20 + " **问题回答** " + "="*20 + f"\n\n{final_answer}"

        except (ValueError, ImportError, RuntimeError) as e:
            return f"处理过程中发生错误: {e}"


# --- 脚本主入口 ---
if __name__ == "__main__":
    # --- 配置区域 ---
    # 本地大模型路径
    llm_model_path = r"/DATA/f60055380/Qwen3_8B_TL/Model/qwen3-8b-hf"
    
    # 开源Embedding模型 (用于RAG)
    embedding_model = "BAAI/bge-large-zh-v1.5" # 这是一个强大的中文embedding模型
    
    # 待处理的文件路径
    file_to_process = r"/DATA/f60055380/Qwen3_8B_TL/data/xlsx/问题集1200.xlsx"
    
    # 你想问的问题
    question = "总结不同维度不同分数的规律。"
    
    # --- 执行 ---
    start_time = time.time()
    try:
        # 初始化高级RAG总结器
        summarizer = AdvancedRAGSummarizer(
            model_name_or_path=llm_model_path,
            embedding_model_name=embedding_model,
            enable_thinking=True  # 启用或禁用“思考模式”
        )
        
        # 执行处理流程
        final_output = summarizer.process_document(
            file_to_process,
            question,
            chunk_size_chars=1000,  # RAG分块大小，通常小一些（500-1000字符）效果更好
            sheet_name="Sheet1" if file_to_process.endswith(('.xlsx', '.xls')) else None
        )
        
        print("\n\n" + "="*25 + " **最终输出结果** " + "="*25)
        print(final_output)
        print("="*70)
        
    except Exception as e:
        print(f"\n处理过程中发生致命错误: {e}")
    
    end_time = time.time()
    print(f"整个任务总用时: {end_time - start_time:.2f} 秒")